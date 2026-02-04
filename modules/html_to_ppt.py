from html.parser import HTMLParser
import io
import tempfile
import subprocess
import sys

import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from playwright.sync_api import sync_playwright
from playwright.sync_api import Error as PlaywrightError


class HTMLTextExtractor(HTMLParser):
    def __init__(self):
        super().__init__()
        self._parts = []

    def handle_data(self, data):
        if data:
            self._parts.append(data)

    def get_text(self):
        return " ".join(part.strip() for part in self._parts if part.strip())


def html_to_text(html):
    parser = HTMLTextExtractor()
    parser.feed(html)
    return parser.get_text() or " "


def render_html_to_png(html, width=1280, height=720):
    with sync_playwright() as playwright:
        browser = playwright.chromium.launch(
            args=["--disable-dev-shm-usage", "--no-sandbox"]
        )
        page = browser.new_page(viewport={"width": width, "height": height})
        page.set_content(html, wait_until="networkidle")
        png_bytes = page.screenshot(full_page=True)
        browser.close()
    return png_bytes


def build_pptx(slide_html_list, include_text_layer, render_mode):
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank_layout = presentation.slide_layouts[6]

    render_as_image = render_mode == "Render HTML as image (requires Playwright Chromium)"

    with tempfile.TemporaryDirectory() as tmp_dir:
        for index, html in enumerate(slide_html_list):
            slide = presentation.slides.add_slide(blank_layout)
            if render_as_image:
                png_bytes = render_html_to_png(html)
                tmp_path = f"{tmp_dir}/slide_{index + 1}.png"
                with open(tmp_path, "wb") as tmp_file:
                    tmp_file.write(png_bytes)
                slide.shapes.add_picture(
                    tmp_path,
                    left=Inches(0),
                    top=Inches(0),
                    width=presentation.slide_width,
                    height=presentation.slide_height,
                )

            if include_text_layer or not render_as_image:
                textbox = slide.shapes.add_textbox(
                    left=Inches(0.6),
                    top=Inches(0.6),
                    width=Inches(12.0),
                    height=Inches(6.3),
                )
                text_frame = textbox.text_frame
                text_frame.word_wrap = True
                paragraph = text_frame.paragraphs[0]
                paragraph.text = html_to_text(html)
                paragraph.font.size = Pt(18)

    buffer = io.BytesIO()
    presentation.save(buffer)
    buffer.seek(0)
    return buffer


def html_to_ppt_ui():
    st.subheader("HTML to PPT")
    st.caption(
        "Paste HTML for each slide, then generate a PPTX rendered from the HTML layout."
    )
    st.info(
        "For visual fidelity, install Playwright Chromium: `playwright install chromium`."
    )
    if st.button("Install Playwright Chromium"):
        with st.spinner("Installing Chromium..."):
            try:
                result = subprocess.run(
                    [sys.executable, "-m", "playwright", "install", "chromium"],
                    check=True,
                    capture_output=True,
                    text=True,
                )
                st.success("Chromium installed successfully.")
                if result.stdout:
                    st.code(result.stdout)
            except subprocess.CalledProcessError as exc:
                st.warning("Standard install failed. Trying `--with-deps`.")
                if exc.stdout:
                    st.code(exc.stdout)
                if exc.stderr:
                    st.code(exc.stderr)
                try:
                    result = subprocess.run(
                        [sys.executable, "-m", "playwright", "install", "--with-deps", "chromium"],
                        check=True,
                        capture_output=True,
                        text=True,
                    )
                    st.success("Chromium installed successfully with dependencies.")
                    if result.stdout:
                        st.code(result.stdout)
                except subprocess.CalledProcessError as deps_exc:
                    st.error("Failed to install Chromium (with deps).")
                    if deps_exc.stdout:
                        st.code(deps_exc.stdout)
                    if deps_exc.stderr:
                        st.code(deps_exc.stderr)

    if "html_slide_count" not in st.session_state:
        st.session_state.html_slide_count = 4

    if "html_slide_inputs" not in st.session_state:
        st.session_state.html_slide_inputs = [""] * st.session_state.html_slide_count

    if st.button("➕ Add Slide"):
        st.session_state.html_slide_count += 1
        st.session_state.html_slide_inputs.append("")

    slide_html_list = []
    for index in range(st.session_state.html_slide_count):
        value = st.text_area(
            f"Slide {index + 1} HTML",
            value=st.session_state.html_slide_inputs[index],
            height=160,
            key=f"html_slide_{index}",
        )
        slide_html_list.append(value)
    st.session_state.html_slide_inputs = slide_html_list

    render_mode = st.radio(
        "Slide rendering mode",
        [
            "Render HTML as image (requires Playwright Chromium)",
            "Text-only (fast, editable, no browser dependency)",
        ],
        index=0,
    )
    include_text_layer = st.checkbox(
        "Add editable text layer (from HTML text)",
        value=True,
        disabled=render_mode != "Render HTML as image (requires Playwright Chromium)",
    )

    if st.button("Generate PPT"):
        if not any(html.strip() for html in slide_html_list):
            st.warning("Please enter HTML for at least one slide.")
            return

        try:
            pptx_buffer = build_pptx(slide_html_list, include_text_layer, render_mode)
        except PlaywrightError as exc:
            if render_mode == "Render HTML as image (requires Playwright Chromium)":
                st.error(
                    "Playwright Chromium is not available or missing dependencies. "
                    "Click 'Install Playwright Chromium', then restart the app if needed, "
                    "or use text-only mode."
                )
                pptx_buffer = build_pptx(
                    slide_html_list,
                    include_text_layer=False,
                    render_mode="Text-only (fast, editable, no browser dependency)",
                )
            else:
                raise exc
        st.success("Your PPT is ready to download.")
        st.download_button(
            "Download PPT",
            data=pptx_buffer,
            file_name="html_slides.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
